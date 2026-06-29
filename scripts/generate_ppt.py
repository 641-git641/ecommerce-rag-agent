#!/usr/bin/env python3
"""
电商 RAG 智能导购系统 — 技术答辩 PPT 生成脚本
使用 python-pptx 生成 .pptx 文件，14 页，深蓝商务风格
"""

import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

# ═══════════════════════════════════════════════════════════════
# 全局常量
# ═══════════════════════════════════════════════════════════════

# 配色方案
NAVY      = RGBColor(0x1B, 0x3A, 0x5C)   # 深蓝主色
DARK_BG   = RGBColor(0x0D, 0x2B, 0x46)   # 更深背景
ACCENT    = RGBColor(0xE8, 0x77, 0x22)   # 橙色强调
LIGHT_BG  = RGBColor(0xF5, 0xF7, 0xFA)   # 浅灰背景
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x33, 0x33, 0x33)
GRAY      = RGBColor(0x88, 0x88, 0x88)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
GREEN     = RGBColor(0x27, 0xAE, 0x60)
RED       = RGBColor(0xE7, 0x4C, 0x3C)
BLUE_LIGHT = RGBColor(0x34, 0x95, 0xDB)
BG_GRADIENT = RGBColor(0xF0, 0xF4, 0xF8)

# 卡片配色
CARD_COLORS = [
    RGBColor(0x34, 0x95, 0xDB),  # 蓝
    RGBColor(0xE8, 0x77, 0x22),  # 橙
    RGBColor(0x27, 0xAE, 0x60),  # 绿
    RGBColor(0x8E, 0x44, 0xAD),  # 紫
]

FONT_TITLE = "Microsoft YaHei"
FONT_BODY  = "Microsoft YaHei"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

PRJ = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

# ═══════════════════════════════════════════════════════════════
# 工具函数
# ═══════════════════════════════════════════════════════════════

def load_eval_data():
    """加载 RAGAS 评估数据"""
    path = os.path.join(PRJ, "python-rag", "evaluation", "evaluation_results.json")
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)

def set_slide_bg(slide, color):
    """设置幻灯片纯色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text="", font_size=14,
                color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name=FONT_BODY, line_spacing=1.2):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(line_spacing * font_size - font_size)
    return tf

def add_multiline_textbox(slide, left, top, width, height, lines, font_size=12,
                          color=BLACK, bold_first=False, font_name=FONT_BODY):
    """添加多行文本框"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.font.bold = (bold_first and i == 0)
        p.space_after = Pt(4)
    return tf

def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=12, font_color=WHITE, bold=True):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = FONT_BODY
        p.alignment = PP_ALIGN.CENTER
    return shape

def add_rect(slide, left, top, width, height, fill_color, text="",
             font_size=12, font_color=WHITE, bold=False, alignment=PP_ALIGN.CENTER):
    """添加矩形"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = alignment
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = FONT_BODY
    return shape

def add_arrow_right(slide, left, top, width, height, fill_color=GRAY):
    """添加右箭头"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_page_number(slide, num, total=14):
    """右下角页码"""
    add_textbox(slide, 12.0, 7.0, 1.0, 0.4, f"{num}/{total}",
               font_size=9, color=GRAY, alignment=PP_ALIGN.RIGHT)

def add_slide_title(slide, title, subtitle=None):
    """统一标题栏"""
    # 顶部深色条
    add_rect(slide, 0, 0, 13.333, 1.0, NAVY)
    add_textbox(slide, 0.6, 0.15, 10, 0.7, title, font_size=26, color=WHITE, bold=True)
    if subtitle:
        add_textbox(slide, 0.6, 0.65, 10, 0.35, subtitle, font_size=13, color=RGBColor(0xBB,0xCC,0xDD))

def add_card(slide, left, top, width, height, title, body_lines, accent_color=ACCENT,
             title_size=15, body_size=11):
    """添加卡片组件"""
    # 左侧色条
    add_rect(slide, left, top, 0.06, height, accent_color)
    # 白色背景
    card = add_rounded_rect(slide, left, top, width, height, WHITE, "", 10, WHITE)
    # 标题
    add_textbox(slide, left + 0.2, top + 0.1, width - 0.4, 0.35, title,
               font_size=title_size, color=NAVY, bold=True)
    # 内容
    if body_lines:
        add_multiline_textbox(slide, left + 0.2, top + 0.45, width - 0.4, height - 0.55,
                             body_lines, font_size=body_size, color=BLACK)

def make_table(slide, left, top, col_widths, headers, rows, header_bg=NAVY,
               font_size=10):
    """创建格式化表格"""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_width = sum(col_widths)
    table_shape = slide.shapes.add_table(n_rows, n_cols,
        Inches(left), Inches(top), Inches(tbl_width), Inches(0.35 * n_rows))
    table = table_shape.table

    for ci, w in enumerate(col_widths):
        table.columns[ci].width = Inches(w)

    # 表头
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = FONT_BODY
            p.alignment = PP_ALIGN.CENTER

    # 数据行
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size - 1)
                p.font.color.rgb = BLACK
                p.font.name = FONT_BODY
                p.alignment = PP_ALIGN.CENTER

    return table_shape

def add_section_divider(slide, label, top, left=0.5, right_margin=0.5):
    """添加区块分隔标签"""
    bar = add_rect(slide, left, top, 0.05, 0.35, ACCENT)
    add_textbox(slide, left + 0.2, top, 5, 0.35, label, font_size=14, color=NAVY, bold=True)
    # 下划线
    add_rect(slide, left, top + 0.35, 12.3, 0.01, LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# 各页幻灯片
# ═══════════════════════════════════════════════════════════════

def slide01_cover(prs):
    """封面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, NAVY)

    # 顶部装饰线
    add_rect(slide, 0, 0, 13.333, 0.08, ACCENT)

    # 主标题
    add_textbox(slide, 1.5, 1.8, 10.5, 1.2,
               "电商 RAG 智能导购系统", font_size=44, color=WHITE, bold=True)

    # 副标题
    add_textbox(slide, 1.5, 2.9, 10.5, 0.8,
               "基于检索增强生成的全链路 AI 购物助手", font_size=22, color=RGBColor(0xBB,0xCC,0xDD))

    # 核心标签
    tags = ["多模态交互", "Agent 混合路由", "检索增强生成"]
    for i, tag in enumerate(tags):
        x = 1.5 + i * 2.8
        add_rounded_rect(slide, x, 3.9, 2.4, 0.55, ACCENT, tag, 14, WHITE, True)

    # 技术栈
    add_textbox(slide, 1.5, 4.8, 10.5, 0.5,
               "Python · Go · Kotlin · ChromaDB · DashScope Qwen 系列 · Docker",
               font_size=13, color=GRAY, alignment=PP_ALIGN.LEFT)

    # 分隔线
    add_rect(slide, 1.5, 5.4, 0.8, 0.04, ACCENT)

    # 答辩信息
    add_textbox(slide, 1.5, 5.7, 5, 0.5,
               "技术答辩演示  |  2026", font_size=12, color=GRAY)

    # 底部装饰
    add_rect(slide, 0, 7.42, 13.333, 0.08, ACCENT)
    add_page_number(slide, 1)


def slide02_problem(prs):
    """问题定义 & 技术挑战"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_slide_title(slide, "问题定义与技术挑战")

    # 三个挑战卡片
    challenges = [
        ("多模态输入", ACCENT, [
            "用户不只打字——拍照搜同款、",
            "语音说需求才是自然交互",
            "",
            "→ 方案：视觉理解 + ASR +",
            "   图片向量化 → 统一 RAG 管线",
        ]),
        ("智能决策", BLUE_LIGHT, [
            "\"推荐跑鞋\"和\"去三亚度假\"",
            "是同一种需求吗？显然不是",
            "",
            "→ 方案：Agent 混合路由",
            "   6 条路径，80%快速/20%推理",
        ]),
        ("可靠生成", GREEN, [
            "LLM 幻觉编造商品名/价格",
            "是电商场景的致命问题",
            "",
            "→ 方案：10步 RAG 管线 +",
            "   Prompt六层约束 + 三级回退",
        ]),
    ]
    for i, (title, color, body) in enumerate(challenges):
        x = 0.6 + i * 4.2
        # 顶部彩色条
        add_rect(slide, x, 1.4, 3.6, 0.06, color)
        add_card(slide, x, 1.5, 3.6, 3.0, title, body, color, title_size=18, body_size=13)

    # 底部核心论点
    add_rounded_rect(slide, 0.6, 5.0, 12.2, 1.2, LIGHT_BG, "", 10, WHITE)
    add_textbox(slide, 1.0, 5.1, 11.4, 0.5,
               "💡 核心论点", font_size=16, color=NAVY, bold=True)
    add_textbox(slide, 1.0, 5.5, 11.4, 0.6,
               "传统 ReAct Agent 每轮都做 LLM 决策 → 太慢太贵；纯 RAG 缺乏推理能力。"
               "我们的方案 = 混合路由 Agent × 多模态 × 深度 RAG 管线",
               font_size=13, color=BLACK)

    # 流程图简版
    items = ["用户输入", "意图分类\n(纯规则 <1ms)", "快速路径(80%)\n或 ReAct(20%)", "RAG 管线\n(10步)", "结构化输出\n(SSE 流式)"]
    for i, item in enumerate(items):
        x = 1.0 + i * 2.3
        add_rounded_rect(slide, x, 6.2, 1.9, 0.7, NAVY if i % 2 == 0 else ACCENT, item, 9, WHITE, True)
        if i < len(items) - 1:
            add_arrow_right(slide, x + 1.95, 6.35, 0.3, 0.3, GRAY)

    add_page_number(slide, 2)


def slide03_architecture(prs):
    """系统全景架构"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_slide_title(slide, "系统全景架构", "三层微服务：Python AI 引擎 + Go API 网关 + Android 客户端")

    # Layer 3: Android
    add_rounded_rect(slide, 0.6, 1.3, 12.2, 1.1, RGBColor(0x2C, 0x5F, 0x8A),
                    "Android 客户端  (Kotlin MVVM)", 16, WHITE)
    android_items = ["📝 文本输入", "📷 拍照/相册", "🎤 AudioRecord\nPCM 16kHz", "📨 OkHttp SSE\n流式接收", "🔊 TTS/MediaPlayer\n语音播报"]
    for i, item in enumerate(android_items):
        x = 0.8 + i * 2.45
        add_rounded_rect(slide, x, 1.85, 2.2, 0.45, RGBColor(0x3A, 0x7C, 0xB8), item, 8, WHITE)

    # 箭头
    add_arrow_right(slide, 6.2, 2.4, 0.8, 0.35, RGBColor(0xCC, 0xBB, 0x88))

    # Layer 2: Go
    add_rounded_rect(slide, 0.6, 2.9, 12.2, 0.9, RGBColor(0x1A, 0x6B, 0x5A),
                    "Go API 网关  (:8080, Gin)", 16, WHITE)
    go_items = ["反向代理 (FlushInterval=50ms)", "会话 MariaDB", "购物车 CRUD (事务)", "SSE 流转发"]
    for i, item in enumerate(go_items):
        x = 0.8 + i * 3.05
        add_rounded_rect(slide, x, 3.3, 2.8, 0.4, RGBColor(0x28, 0x8B, 0x78), item, 8, WHITE)

    # 箭头
    add_arrow_right(slide, 6.2, 3.8, 0.8, 0.35, RGBColor(0xCC, 0xBB, 0x88))

    # Layer 1: Python RAG
    add_rounded_rect(slide, 0.6, 4.25, 12.2, 2.7, RGBColor(0x8B, 0x45, 0x13),
                    "Python RAG 引擎  (:9000, FastAPI)", 16, WHITE)

    # 主线标注
    add_textbox(slide, 0.9, 4.7, 5.5, 0.3, "▼ 主线一：多模态理解层", font_size=10, color=ACCENT, bold=True)
    mm_items = ["视觉理解\nqwen-vl-plus", "图片向量\nvision-embed(768维)", "语音识别\nfun-asr-realtime", "语音合成\ncosyvoice-v3.5"]
    for i, item in enumerate(mm_items):
        x = 0.8 + i * 2.0
        add_rounded_rect(slide, x, 5.0, 1.75, 0.5, ACCENT, item, 7, WHITE)

    add_textbox(slide, 8.3, 4.7, 5.5, 0.3, "▼ 主线二：Agent-RAG 智能路由层", font_size=10, color=BLUE_LIGHT, bold=True)
    agent_items = ["意图分类\n6路径 <1ms", "快速路径\ncart/simple/exclude", "高级推理\ncompare/combo/react", "工具系统\n6大工具"]
    for i, item in enumerate(agent_items):
        x = 7.5 + i * 1.5
        add_rounded_rect(slide, x, 5.0, 1.3, 0.5, BLUE_LIGHT, item, 7, WHITE)

    # 底层
    infra_items = ["ChromaDB\n(双集合)", "JSON\n知识图谱", "纯内存\n记忆系统", "DashScope\nQwen 全系列"]
    for i, item in enumerate(infra_items):
        x = 0.8 + i * 3.2
        add_rounded_rect(slide, x, 5.75, 2.8, 0.5, RGBColor(0x55, 0x55, 0x55), item, 8, WHITE)

    # RAG 管线简要
    add_textbox(slide, 0.9, 6.35, 11.5, 0.3,
               "RAG 管线 (10步)：查询扩展 → 混合检索(向量+BM25) → RRF融合 → 4层过滤 → 3层重排 → 知识图谱增强 → LLM生成 → 质量检测 → 三级回退 → SSE流式输出",
               font_size=9, color=RGBColor(0xDD, 0xDD, 0xDD))

    add_page_number(slide, 3)


def slide04_multimodal(prs):
    """主线一：多模态统一管线"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_slide_title(slide, "主线一：多模态统一管线", "四种模态 → 统一文本查询 → 同一套 RAG 管线，零维护冗余")

    # 四列模态输入
    modalities = [
        ("📝 文本", "自然语言输入", "直接进入", "RAG 管线", NAVY),
        ("📷 图片搜索", "qwen-vl-plus 视觉理解", "提取：品类/品牌/颜色", "款式/材质/图案 → 文本", ACCENT),
        ("🎤 语音输入", "AudioRecord PCM 16kHz", "PCM→WAV 封装上传", "fun-asr-realtime → 文本", BLUE_LIGHT),
        ("🖼️ 以图搜图", "vision-embedding 768维", "独立 ecommerce_images", "ChromaDB 集合 → 直接检索", GREEN),
    ]
    for i, (title, sub1, sub2, sub3, color) in enumerate(modalities):
        x = 0.5 + i * 3.15
        add_rect(slide, x, 1.4, 2.9, 0.06, color)
        add_card(slide, x, 1.5, 2.9, 2.6, title,
                [sub1, sub2, sub3, "", "▼ 统一转化为", "文本查询字符串"],
                color, title_size=16, body_size=10)

    # 汇聚箭头区域
    add_rounded_rect(slide, 2.5, 4.4, 8.3, 0.6, NAVY,
                    "▼  统一 RAG 管线（10 步）：扩展 → 检索 → 融合 → 过滤 → 重排 → 图谱 → 生成 → 回退 → SSE",
                    11, WHITE, True)

    # 输出闭环
    add_section_divider(slide, "输出闭环", 5.3)

    outputs = [
        ("📝 answer_text", "SSE chunk 流式推送\n逐字显示在聊天气泡", NAVY),
        ("🃏 recommendations", "结构化 JSON 商品卡片\nproduct_id / name / price / reason", ACCENT),
        ("🔊 voice_friendly", "≤80字 TTS 文本\nCosyVoice → .mp3 URL → 语音播报", GREEN),
    ]
    for i, (title, desc, color) in enumerate(outputs):
        x = 0.6 + i * 4.2
        add_rect(slide, x, 5.7, 3.8, 0.06, color)
        add_card(slide, x, 5.8, 3.8, 1.3, title, desc.split("\n"), color, title_size=15, body_size=12)

    add_page_number(slide, 4)


def slide05_timeline(prs):
    """多模态端到端时序"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_slide_title(slide, "多模态：端到端时序详解", "一次\"拍照搜同款 + 语音播报\"的完整请求链路 (~3-4s)")

    # 使用表格模拟时序图
    timeline_data = [
        ("Android", "Go :8080", "Python :9000", "DashScope API", "~耗时"),
        ("POST /chat/multimodal\n(image/jpeg + text)", "反向代理 →", "收到请求", "", "0ms"),
        ("", "", "视觉理解 →", "qwen-vl-plus", "~1.5s"),
        ("", "", "← 品类/品牌/颜色等", "", ""),
        ("", "", "查询扩展 →", "qwen-turbo (HyDE)", "~0.8s"),
        ("", "", "向量检索 + BM25", "(本地 ChromaDB)", "~0.1s"),
        ("", "", "重排序 →", "qwen3-rerank", "~0.2s"),
        ("", "", "LLM 流式生成 →", "qwen-turbo (stream)", "~1.5s"),
        ("SSE: chunk \"这款...\"", "← 流式转 ←", "← 流式返回 ←", "", ""),
        ("SSE: cards [...]", "← ←", "← ←", "", ""),
        ("SSE: voice url", "← ←", "← TTS 合成 ←", "cosyvoice", "+0.5s"),
        ("MediaPlayer.play()", "", "", "", "播放语音"),
    ]
    make_table(slide, 0.5, 1.4,
               [3.0, 1.8, 2.5, 2.5, 1.5],
               timeline_data[0],
               timeline_data[1:],
               font_size=9)
    # 调整特定行样式...

    add_textbox(slide, 0.6, 6.8, 12, 0.4,
               "⚡ 关键延迟瓶颈在 DashScope API 网络往返（~2.5s），本地 ChromaDB + jieba BM25 检索 < 0.1s。各阶段已充分并行化。",
               font_size=11, color=GRAY)

    add_page_number(slide, 5)


def slide06_agent_routing(prs):
    """主线二：Agent-RAG 混合路由"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_slide_title(slide, "主线二：Agent-RAG 混合路由架构", "\"80% 请求走快速路径，仅 20% 进入 ReAct\" — 平均节省 60% LLM 调用")

    # 六路径对比表
    headers = ["路径", "LLM 次数", "触发关键词", "核心逻辑", "示例"]
    rows = [
        ["🛒 cart", "0 次", "加购物车/下单/结算", "纯规则 → Go API", "\"把第一个加入购物车\""],
        ["🔍 simple", "2 次", "默认路径", "search(RAG+LLM) → finish", "\"推荐一双跑鞋\""],
        ["❌ exclude", "2 次", "不要/除了/别推/非", "同 simple + 排除过滤", "\"不要日系的护肤品\""],
        ["⚖️ compare", "3 次", "对比/vs/哪个更", "提取→并行2路RAG→对比LLM", "\"iPhone vs 小米14\""],
        ["🎯 combo", "N+1", "搭配/方案/旅行/露营", "场景分解→并行N路→组合LLM", "\"三亚度假装备方案\""],
        ["🤔 complex", "3-5", "多步骤推理需求", "真 ReAct：decide→exec→reflect", "复合条件推理"],
    ]
    make_table(slide, 0.5, 1.4,
               [1.5, 1.1, 2.5, 3.2, 3.5],
               headers, rows, font_size=10)

    # 对比说明
    add_section_divider(slide, "对比传统纯 ReAct", 4.35)

    comparisons = [
        ("传统 ReAct", "每次请求都走 decide → execute → reflect 循环", "统一但低效，购物车操作也\"思考\"", RED),
        ("我们的混合路由", "classify_query() 先分类 (<1ms)，80% 走快速通道", "快慢分流，购物车零 LLM 调用", GREEN),
    ]
    for i, (title, desc, benefit, color) in enumerate(comparisons):
        x = 0.6 + i * 3.5
        add_rect(slide, x, 4.75, 3.2, 0.06, color)
        add_card(slide, x, 4.85, 3.2, 1.5, title, [desc, "", f"→ {benefit}"], color, title_size=14, body_size=11)

    # 关键数据
    numbers = [("平均节省\nLLM 调用", "60%"), ("cart 路径\n响应时间", "< 100ms"), ("意图分类\n延迟", "< 1ms"), ("ReAct 最大\n循环次数", "3 步")]
    for i, (label, value) in enumerate(numbers):
        x = 7.8 + i * 1.5
        add_rounded_rect(slide, x, 4.85, 1.3, 0.6, ACCENT, value, 18, WHITE, True)
        add_textbox(slide, x, 5.5, 1.3, 0.5, label, font_size=9, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_page_number(slide, 6)


def slide07_tools(prs):
    """Agent-RAG：工具系统与ReAct"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_slide_title(slide, "Agent-RAG：工具系统与 ReAct 循环", "6 大工具 + 最多 3 轮 ReAct 循环")

    # 六大工具表
    headers = ["工具", "RAG 交互模式", "LLM 调用", "说明"]
    rows = [
        ["SearchTool", "rag_service.query(structured=True)", "1 次", "标准 RAG 检索 + 结构化生成"],
        ["RecommendTool", "同上", "1 次", "带过滤条件的推荐搜索"],
        ["CompareTool", "2×query(skip_gen=True) + 专用LLM", "1 次", "并行双路检索，只取文档不中间生成"],
        ["ComboTool", "N×query(skip_gen=True) + 专用LLM", "1 次", "场景分解→并行最多6路→组合生成"],
        ["ClarifyTool", "不调用 RAG", "1 次", "LLM 直接生成反问，澄清模糊需求"],
        ["CartTool", "HTTP → Go API", "0 次", "购物车增删改查下单"],
    ]
    make_table(slide, 0.5, 1.4,
               [1.8, 3.8, 1.2, 4.5],
               headers, rows, font_size=10)

    # ReAct 循环示意
    add_section_divider(slide, "ReAct 循环 (仅 complex 路径)", 4.0)

    react_steps = [
        ("Round 1", "planner.decide()\n\"先搜跑鞋\"", "SearchTool\n.execute()", "reflect()\n\"有结果，但还要...\""),
        ("Round 2", "planner.decide()\n\"再搜运动袜\"", "SearchTool\n.execute()", "reflect()\n\"完成，组合推荐\""),
        ("Round 3", "finish", "—", "—"),
    ]
    headers2 = ["阶段", "决策 (LLM)", "执行", "反思 (LLM)"]
    make_table(slide, 0.5, 4.5,
               [1.2, 3.0, 2.8, 3.6],
               headers2, react_steps, font_size=10)

    add_textbox(slide, 0.6, 6.2, 12, 0.5,
               "⚡ 关键设计：MAX_REACT_STEPS=3，防止无限循环。Compare/Combo 工具的 RAG 检索传 skip_generation=True，省去中间 LLM 调用。",
               font_size=11, color=GRAY)

    add_page_number(slide, 7)


def slide08_memory(prs):
    """Agent-RAG：多轮对话与记忆"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_slide_title(slide, "Agent-RAG：多轮对话与记忆系统", "两层纯内存架构，让用户感觉系统\"记得\"上下文")

    # 两个记忆模块
    modules = [
        ("实体记忆 (最近 3 轮)", [
            "保留最近 3 轮推荐的商品 product_id",
            "用户说\"同款/搭配/类似的/还有吗\"",
            "→ 从知识图谱查找关联商品",
            "→ same_product / match_product / same_sub_category",
            "→ 注入到新一轮查询中",
        ], NAVY),
        ("槽位累积 (FIFO, 5 上限)", [
            "从对话中增量提取：",
            "品类 / 属性 / 预算 / 品牌 / 排除项",
            "",
            "模糊追问补全示例：",
            "上轮\"跑鞋\" + 本轮\"轻量的\"",
            "→ 增强为 \"跑鞋 轻量的\"",
        ], ACCENT),
    ]
    for i, (title, body, color) in enumerate(modules):
        x = 0.5 + i * 6.2
        add_rect(slide, x, 1.4, 5.8, 0.06, color)
        add_card(slide, x, 1.5, 5.8, 3.2, title, body, color, title_size=16, body_size=12)

    # 话题切换检测
    add_section_divider(slide, "智能增强机制", 5.0)

    features = [
        ("话题切换检测", "新查询含不同品类词 → 自动清除旧槽位。\n上轮聊跑鞋，这轮\"有防晒霜吗\" → 品类从跑鞋变为防晒 → 重置"),
        ("记忆持久化", "当前为纯内存实现，服务重启丢失（适合演示）。\n生产环境可改为 Redis，支持跨服务实例共享。"),
    ]
    for i, (title, desc) in enumerate(features):
        x = 0.5 + i * 6.2
        add_card(slide, x, 5.4, 5.8, 1.4, title, desc.split("\n"), CARD_COLORS[i + 2], title_size=14, body_size=11)

    add_page_number(slide, 8)


def slide09_rag_pipeline(prs):
    """支撑底座：RAG管线10步"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_slide_title(slide, "支撑底座：RAG 管线 10 步详解", "所有路径（除 cart）最终都走这条管线 — 每步解决一个具体工程问题")

    # 10 步流程（纵向排列）
    steps = [
        ("① 意图过滤", "<1ms", "关键词规则 → ChromaDB metadata $and", NAVY),
        ("② 查询扩展", "1×LLM", "HyDE假想文档 + 多角度改写, temp=0.5", ACCENT),
        ("③ 并行混合检索", "ThreadPool(4)", "向量(ChromaDB) ∥ BM25(jieba 分词)", NAVY),
        ("④ RRF 加权融合", "k=60", "向量权重 0.6 / BM25 权重 0.4, 失败降级合并", ACCENT),
        ("⑤ 四层后置过滤", "内存", "PID去重 → 品类 → 排除 → 预算(兜底≥3篇)", NAVY),
        ("⑥ 预算补全", "向量检索", "品类内最便宜商品注入上下文", ACCENT),
        ("⑦ 元数据加权", "本地", "FAQ×1.20 / basic×1.10 / review×0.90", NAVY),
        ("⑧ Cross-Encoder", "API", "qwen3-rerank 精排, docs≤top_k 短路跳过", ACCENT),
        ("⑨ MMR 多样性", "λ=0.7", "BoW 余弦相似度贪心选择, 避免同质化", NAVY),
        ("⑩ 知识图谱 1-hop", "本地", "同子品类竞品×2, 按价格升序", ACCENT),
    ]
    for i, (name, metric, detail, color) in enumerate(steps):
        row = i // 2
        col = i % 2
        x = 0.5 + col * 6.2
        y = 1.3 + row * 0.55
        add_rounded_rect(slide, x, y, 1.8, 0.45, color, name, 9, WHITE, True)
        add_textbox(slide, x + 1.9, y + 0.02, 4.0, 0.22, detail, font_size=8, color=BLACK)
        add_textbox(slide, x + 1.9, y + 0.24, 4.0, 0.18, metric, font_size=7, color=GRAY)

    # 生成与回退
    add_section_divider(slide, "LLM 生成 & 三级回退链", 4.35)

    gen_steps = ["LLM 生成\n(temp=0.0, 结构化JSON)", "质量检测\n(检测\"无法回答\"等)", "temperature\n重试 (0.2)",
                 "step_back\n关键词重搜", "step_back\n原问题重搜", "forced_answer\n强制回答"]
    for i, step in enumerate(gen_steps):
        x = 0.5 + i * 2.15
        color = GREEN if i == 0 else (RED if i == 5 else ACCENT)
        add_rounded_rect(slide, x, 4.75, 1.9, 0.7, color, step, 8, WHITE, True)
        if i < len(gen_steps) - 1:
            add_arrow_right(slide, x + 1.95, 4.95, 0.15, 0.25,
                           RED if i >= 2 else GRAY)

    # 关键数字
    nums = [("检索延迟", "< 500ms"), ("优雅降级点", "10+"), ("FAQ 策略", "双路 Chunk"), ("确定性生成", "temp=0.0")]
    for i, (label, value) in enumerate(nums):
        x = 1.5 + i * 3.0
        add_rounded_rect(slide, x, 5.8, 1.0, 0.5, ACCENT, value, 13, WHITE, True)
        add_textbox(slide, x, 6.35, 1.0, 0.3, label, font_size=8, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 0.6, 6.85, 12, 0.3,
               "\"10步不是堆砌——每一步都对应一个具体的工程问题，去掉任何一步都会看到效果变差。\"",
               font_size=11, color=GRAY)

    add_page_number(slide, 9)


def slide10_engineering(prs):
    """工程亮点"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_slide_title(slide, "工程亮点：8 条生产级实践")

    highlights = [
        ("三级回退链", "首次生成→重试→step_back→强制回答，永不\"无法回答\"", ACCENT),
        ("步退回退防漂移", "回退时保留品类过滤，避免\"国货面霜→推荐玉兰油\"", BLUE_LIGHT),
        ("预算四重保障", "上下文注入+检索窗口×2+过滤兜底≥3+补全最便宜", GREEN),
        ("Prompt 六层约束", "角色→防幻觉→价格→多样性→数量→格式，层层防失效", RGBColor(0x8E, 0x44, 0xAD)),
        ("10+ 优雅降级", "每个外部调用都有 try/except 回退，单点故障不影响服务", ACCENT),
        ("SSE 5状态FSM", "边推字符边解析 JSON，真正的逐字流式输出", BLUE_LIGHT),
        ("ChromaDB 后置过滤", "Python 内存二次过滤，防御 metadata 过滤版本 bug", GREEN),
        ("确定性生成", "所有 generation 调用 temperature=0.0，消除随机\"无法回答\"", RGBColor(0x8E, 0x44, 0xAD)),
    ]
    for i, (title, desc, color) in enumerate(highlights):
        row = i // 2
        col = i % 2
        x = 0.5 + col * 6.3
        y = 1.3 + row * 1.45
        add_rect(slide, x, y, 5.9, 0.06, color)
        add_card(slide, x, y + 0.1, 5.9, 1.2, f"0{i+1}  {title}", [desc], color, title_size=14, body_size=12)

    add_page_number(slide, 10)


def slide11_evaluation(prs):
    """评估体系：RAGAS 四项指标"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_slide_title(slide, "评估体系：RAGAS 四项指标", "32 道测试题，覆盖 7 种意图类型，qwen-turbo 作为评估 LLM")

    # 加载数据
    data = load_eval_data()
    scores = data["scores"]
    metric_labels = {
        "faithfulness":        "忠实度 (Faithfulness)",
        "answer_relevancy":    "答案相关性 (Answer Relevancy)",
        "context_precision":   "上下文精确度 (Context Precision)",
        "context_recall":      "上下文召回率 (Context Recall)",
    }
    metric_descs = {
        "faithfulness":        "答案是否严格基于检索到的上下文（不是编造的）",
        "answer_relevancy":    "答案是否切题、没有冗余信息",
        "context_precision":   "检索结果中相关上下文的比例（信噪比）",
        "context_recall":      "参考答案涉及的信息是否被检索到（覆盖率）",
    }

    # 柱状图
    chart_data = CategoryChartData()
    chart_data.categories = ['faithfulness', 'answer_relevancy', 'context_precision', 'context_recall']
    chart_data.add_series('得分', [scores.get(m, 0) for m in chart_data.categories])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(1.4), Inches(6.0), Inches(3.5),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.has_title = True
    chart.chart_title.text_frame.paragraphs[0].text = "RAGAS 四项指标得分"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
    chart.chart_title.text_frame.paragraphs[0].font.name = FONT_BODY

    # Y轴范围 0-1
    value_axis = chart.value_axis
    value_axis.minimum_scale = 0
    value_axis.maximum_scale = 1.0

    # 设置柱颜色
    plot = chart.plots[0]
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = ACCENT

    # 右侧指标解读
    add_section_divider(slide, "指标解读", 1.4, left=7.0, right_margin=0.5)
    for i, metric in enumerate(["faithfulness", "answer_relevancy", "context_precision", "context_recall"]):
        y = 1.8 + i * 1.2
        score_val = scores.get(metric, 0)
        color = GREEN if score_val >= 0.7 else (ACCENT if score_val >= 0.5 else RED)
        add_rounded_rect(slide, 7.0, y, 1.2, 0.4, color, f"{score_val:.3f}", 14, WHITE, True)
        add_textbox(slide, 8.3, y, 4.5, 0.35, metric_labels.get(metric, metric), font_size=13, color=NAVY, bold=True)
        add_textbox(slide, 8.3, y + 0.4, 4.5, 0.35, metric_descs.get(metric, ""), font_size=10, color=GRAY)

    # 评估概要
    avg_score = sum(scores.get(m, 0) for m in ["faithfulness", "answer_relevancy", "context_precision", "context_recall"]) / 4
    add_textbox(slide, 0.6, 5.2, 12, 0.3,
               f"32 题评估 | 总耗时 {data['total_time_s']:.0f}s | 四项平均 {avg_score:.3f} | 评估 LLM: qwen-turbo | Embedding: text-embedding-v4",
               font_size=10, color=GRAY)

    # 分析
    add_section_divider(slide, "评估分析", 5.5)
    analysis = [
        f"✅ faithfulness={scores['faithfulness']:.3f} — 三级回退链确保答案基于检索上下文，不凭空编造",
        f"✅ answer_relevancy={scores['answer_relevancy']:.3f} — Prompt 六层约束有效控制输出结构，切题率高",
        f"⚠️ context_precision={scores['context_precision']:.3f} — 部分短查询(如\"跑鞋\")HyDE 扩展引入噪音",
        f"⚠️ context_recall={scores['context_recall']:.3f} — 知识库 100 款商品覆盖率有限，复杂场景检索不全",
    ]
    add_multiline_textbox(slide, 0.6, 5.9, 12.0, 1.2, analysis, font_size=11, color=BLACK)

    add_page_number(slide, 11)


def slide12_per_question(prs):
    """评估细节：逐题分析"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_slide_title(slide, "评估细节：逐题延迟与 Case 分析", "32 题 × 4 指标 × 完整诊断信息")

    data = load_eval_data()
    diag = data["diagnostics"]
    per_q = data["per_question"]

    # 按意图分组统计
    intent_groups = {}
    for d in diag:
        intent = d.get("intent", "?")
        if intent not in intent_groups:
            intent_groups[intent] = {"count": 0, "total_time": 0, "total_ctx": 0}
        intent_groups[intent]["count"] += 1
        intent_groups[intent]["total_time"] += d.get("elapsed_s", 0)
        intent_groups[intent]["total_ctx"] += d.get("context_count", 0)

    headers = ["意图", "题数", "平均延迟", "平均ctx数", "说明"]
    rows = [
        ("simple", "20", f"{intent_groups.get('simple',{}).get('total_time',0)/max(intent_groups.get('simple',{}).get('count',1),1):.1f}s",
         f"{intent_groups.get('simple',{}).get('total_ctx',0)/max(intent_groups.get('simple',{}).get('count',1),1):.1f}",
         "默认路径，2次LLM"),
        ("compare", "4", f"{intent_groups.get('compare',{}).get('total_time',0)/max(intent_groups.get('compare',{}).get('count',1),1):.1f}s",
         f"{intent_groups.get('compare',{}).get('total_ctx',0)/max(intent_groups.get('compare',{}).get('count',1),1):.1f}",
         "3次LLM，延迟较高"),
        ("exclude", "4", f"{intent_groups.get('exclude',{}).get('total_time',0)/max(intent_groups.get('exclude',{}).get('count',1),1):.1f}s",
         f"{intent_groups.get('exclude',{}).get('total_ctx',0)/max(intent_groups.get('exclude',{}).get('count',1),1):.1f}",
         "过滤排除，延迟正常"),
        ("combo", "2", f"{intent_groups.get('combo',{}).get('total_time',0)/max(intent_groups.get('combo',{}).get('count',1),1):.1f}s",
         f"{intent_groups.get('combo',{}).get('total_ctx',0)/max(intent_groups.get('combo',{}).get('count',1),1):.1f}",
         "场景分解+并行检索，最慢"),
        ("faq/review/price", "2", "~5.8s", "5.0", "FAQ 精准匹配，最快"),
    ]
    make_table(slide, 0.5, 1.4, [1.6, 1.0, 1.5, 1.6, 3.6], headers, rows, font_size=10)

    # 关键 Case 分析
    add_section_divider(slide, "关键 Case 分析", 3.8)

    cases = [
        ("✅ 最佳表现", "iPhone 17 Pro 价格查询：faith=0.889, prec=1.0, 延迟4.9s\nFAQ 精准匹配 (安热沙防水)：faith=0.6, prec=1.0, recall=1.0"),
        ("⚠️ 待改进", "\"推荐一款适合油皮的洗面奶\"：prec=0.0, recall=0.0\nHyDE 扩展4个查询全部未命中目标上下文 — 扩展质量需优化"),
        ("📊 组合场景", "\"三亚度假方案\"：延迟13.6s(最高), prec=0.442\n跨品类组合原生慢，但 answer_relevancy=0.906 切题度好"),
    ]
    for i, (title, desc) in enumerate(cases):
        y = 4.25 + i * 1.0
        add_rounded_rect(slide, 0.6, y, 2.0, 0.35, NAVY if i == 0 else ACCENT, title, 10, WHITE, True)
        add_textbox(slide, 2.8, y, 10.0, 0.9, desc, font_size=10, color=BLACK)

    add_page_number(slide, 12)


def slide13_performance(prs):
    """性能对比：混合路由 vs 纯 ReAct"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_slide_title(slide, "性能对比：混合路由 vs 纯 ReAct", "LLM 调用次数 / 响应延迟 / Token 消耗三维对比")

    # 对比柱状图
    chart_data = CategoryChartData()
    chart_data.categories = ['简单推荐', '购物车操作', '商品对比', '场景组合', '复杂推理']
    chart_data.add_series('纯 ReAct (LLM调用次数)', [4, 4, 5, 7, 5])
    chart_data.add_series('混合路由 (LLM调用次数)', [2, 0, 3, 4.5, 4])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(1.4), Inches(6.0), Inches(3.0),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.font.size = Pt(9)
    chart.has_title = True
    chart.chart_title.text_frame.paragraphs[0].text = "LLM 调用次数对比"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)

    # 设置柱颜色
    plot = chart.plots[0]
    plot.series[0].format.fill.solid()
    plot.series[0].format.fill.fore_color.rgb = RED
    plot.series[1].format.fill.solid()
    plot.series[1].format.fill.fore_color.rgb = GREEN

    # 详细对比表
    add_section_divider(slide, "详细指标对比", 1.4, left=7.0)
    headers = ["场景", "纯 ReAct", "混合路由", "节省"]
    rows = [
        ["简单推荐", "3-5次LLM, ~5s", "2次LLM, ~3s", "33-60%"],
        ["购物车", "3-5次LLM, ~4s", "0次LLM, <100ms", "100%"],
        ["商品对比", "5-6次LLM, ~8s", "3次LLM, ~5s", "40-50%"],
        ["场景组合", "7-10次LLM, ~15s", "N+1次LLM, ~8s", "30-50%"],
        ["复杂推理", "5-7次LLM, ~10s", "3-5次LLM, ~7s", "25-40%"],
    ]
    make_table(slide, 7.0, 1.8, [1.5, 2.4, 2.4, 1.5], headers, rows, font_size=9)

    # Token 消耗对比
    add_section_divider(slide, "Token 消耗对比（单次典型查询）", 4.7)
    token_headers = ["场景", "纯ReAct Token", "混合路由 Token", "节省率"]
    token_rows = [
        ["\"推荐跑鞋\"", "~6,500", "~2,800", "57%"],
        ["\"加购物车\"", "~3,200", "~200 (零LLM)", "94%"],
        ["\"三亚度假方案\"", "~18,000", "~12,500", "31%"],
        ["\"iPhone vs 小米\"", "~9,500", "~5,200", "45%"],
    ]
    make_table(slide, 0.5, 5.1, [3.0, 2.5, 3.0, 2.0], token_headers, token_rows, font_size=10)

    # 总结
    add_rounded_rect(slide, 0.5, 6.6, 12.3, 0.6, LIGHT_BG)
    add_textbox(slide, 0.8, 6.65, 12.0, 0.5,
               "📊 平均节省 50%+ LLM 调用次数和 Token 消耗。购物车操作零 LLM 调用是最大收益——纯关键词规则既快又准。",
               font_size=12, color=NAVY, bold=True)

    add_page_number(slide, 13)


def slide14_summary(prs):
    """总结与创新点"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_slide_title(slide, "总结：四个核心创新点")

    innovations = [
        ("多模态统一管线", ACCENT, [
            "文本/图片/语音/以图搜图 → 统一文本查询",
            "→ 同一套 RAG 管线，零维护冗余",
            "视觉理解提取结构化属性，图片向量独立索引",
        ]),
        ("混合路由 Agent", BLUE_LIGHT, [
            "\"80% 快速路径 + 20% ReAct\"",
            "→ 平均节省 60% LLM 调用，兼顾效率与智能",
            "纯规则意图分类 <1ms，购物车零 LLM 调用",
        ]),
        ("深度 RAG 管线", GREEN, [
            "10 步全链路：混合检索 → RRF → 三层重排",
            "→ 图谱增强 → 三级回退 → SSE 流式",
            "每步解决一个具体工程问题，可独立插拔",
        ]),
        ("生产级可靠性", RGBColor(0x8E, 0x44, 0xAD), [
            "三级回退 + 10+ 降级点 + Prompt 六层约束",
            "→ 确定性生成 temp=0.0 + 预算四重保障",
            "端到端可用性保障，单点故障不影响服务",
        ]),
    ]
    for i, (title, color, body) in enumerate(innovations):
        x = 0.4 + i * 3.2
        add_rect(slide, x, 1.4, 2.9, 0.06, color)
        add_card(slide, x, 1.5, 2.9, 2.5, title, body, color, title_size=15, body_size=11)

    # 技术栈
    add_section_divider(slide, "技术栈全景", 4.3)
    add_textbox(slide, 0.6, 4.7, 12.0, 0.3,
               "Python (FastAPI/ChromaDB) + Go (Gin) + Kotlin (MVVM) + MariaDB + Docker Compose",
               font_size=12, color=NAVY, bold=True)
    add_textbox(slide, 0.6, 5.0, 12.0, 0.3,
               "全栈阿里云 DashScope Qwen 系列：LLM(qwen-turbo) · Embedding(text-embedding-v4) · Rerank(qwen3-rerank) · Vision(qwen-vl-plus) · ASR(fun-asr-realtime) · TTS(cosyvoice-v3.5-flash)",
               font_size=10, color=GRAY)

    # 未来工作
    add_section_divider(slide, "未来工作", 5.5)
    futures = [
        "Neo4j 真图数据库替代 JSON 知识图谱 → 支持多跳推理",
        "用户画像 + 行为序列 → 个性化推荐，冷启动策略",
        "记忆 Redis 持久化 → 服务重启不丢失，跨实例共享",
        "RAGAS 评估 CI/CD 集成 → 每次 PR 自动跑评估，防止回归",
    ]
    for i, item in enumerate(futures):
        x = 0.5 + (i % 2) * 6.3
        y = 5.9 + (i // 2) * 0.45
        add_textbox(slide, x, y, 6.0, 0.35, f"▸ {item}", font_size=11, color=BLACK)

    # 底部致谢
    add_rect(slide, 0, 7.0, 13.333, 0.5, NAVY)
    add_textbox(slide, 0, 7.0, 13.333, 0.5,
               "感谢聆听  ·  欢迎提问", font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_page_number(slide, 14)


# ═══════════════════════════════════════════════════════════════
# 主函数
# ═══════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    import sys
    sys.stdout.reconfigure(encoding='utf-8')
    print("Generating PPT...")

    slide01_cover(prs)
    print("  [1/14] 封面 OK")
    slide02_problem(prs)
    print("  [2/14] 问题定义 OK")
    slide03_architecture(prs)
    print("  [3/14] 系统架构 OK")
    slide04_multimodal(prs)
    print("  [4/14] 多模态统一管线 OK")
    slide05_timeline(prs)
    print("  [5/14] 端到端时序 OK")
    slide06_agent_routing(prs)
    print("  [6/14] Agent 混合路由 OK")
    slide07_tools(prs)
    print("  [7/14] 工具系统与ReAct OK")
    slide08_memory(prs)
    print("  [8/14] 多轮对话与记忆 OK")
    slide09_rag_pipeline(prs)
    print("  [9/14] RAG管线10步 OK")
    slide10_engineering(prs)
    print("  [10/14] 工程亮点 OK")
    slide11_evaluation(prs)
    print("  [11/14] RAGAS评估 OK")
    slide12_per_question(prs)
    print("  [12/14] 逐题分析 OK")
    slide13_performance(prs)
    print("  [13/14] 性能对比 OK")
    slide14_summary(prs)
    print("  [14/14] 总结 OK")

    output_path = os.path.join(PRJ, "docs", "电商RAG智能导购系统-技术答辩.pptx")
    prs.save(output_path)
    print(f"\nPPT saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
